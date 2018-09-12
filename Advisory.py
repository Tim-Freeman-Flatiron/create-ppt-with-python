from pptx import Presentation
from pptx.util import Inches, Pt
from os import getcwd as get_current_directory

class Advisory():
    def __init__(self, students):
        self.advisories = self.make_advisories(students)
        self.advisory_names = self.advisories.keys()
        self.NYU = self.advisories['NYU']

    def make_advisories(self, students):
        advisories = {}
        for stu_id in students:
            student = students[stu_id]
            advisory = student['advisory']
            if not advisories.get(advisory, None):
                advisories[advisory] = []
            advisories[advisory].append(student)
        return advisories

    def get_pw_jumpers(self, students):
        jumpers = list(filter((lambda student: int(student['data'][1:]) > 0), self.get_advisory_data_group(students, 'pw_change')))
        return sorted(jumpers, key=lambda student: int(student['data'][1:]))

    def get_gpa_jumpers(self, students):
        jumpers = list(filter((lambda student: float(student['data'][1:]) > 0), self.get_advisory_data_group(students, 'gpa_change')))
        return sorted(jumpers, key=lambda student: float(student['data'][1:]))

    def get_gpa_leaders(self, students):
        leaders = list(filter((lambda student: float(student['data']) >= 3.0), self.get_advisory_data_group(students, 'gpa')))
        return sorted(leaders, key=lambda student: float(student['data']))

    def get_advisory_data_group(self, students_list, attribute):
        separated_students = []
        for student in students_list:
            temp_obj = {}
            temp_obj['first_name'] = student['first_name']
            temp_obj['last_name'] = student['last_name']
            temp_obj['advisory'] = student['advisory']
            temp_obj['data'] = student[attribute]
            separated_students.append(temp_obj)

        # separated_students = sorted(sortable_data, key=lambda student: student['data'], reverse=True)
        for student in separated_students:
            if 'change' in attribute:
                student['data'] = '+' + str(student['data'])
            else:
                student['data'] = str(student['data'])
        return separated_students

    def add_student_to_slide(self, slide, student, textbox):
        first_name = student['first_name']
        last_name = student['last_name']
        data = student['data']
        textbox.text_frame.text = "{} {} {}".format(first_name, last_name, data)

    def make_ppt(self):
        print('Making PowerPoint...')
        presentation = Presentation()
        blank_slide_layout = presentation.slide_layouts[6]

        for name in self.advisory_names:
            students = self.advisories[name]
            groups = [
                ["Prep Work Jumpers",self.get_pw_jumpers(students)],
                ["GPA Jumpers", self.get_gpa_jumpers(students)],
                ["GPA Leaders", self.get_gpa_leaders(students)]
            ]

            for group in groups:
                group_name = group[0]
                data = group[1]
                slide = presentation.slides.add_slide(blank_slide_layout)
                row = 0
                number_of_students_on_slide = 0
                left_margin = 0.5
                top_margin = 1
                textbox_width = 1
                textbox_height = 1
                title_box = slide.shapes.add_textbox(Inches(left_margin),Inches(top_margin),Inches(textbox_width),Inches(textbox_height))
                title_box.text_frame.text = name + " " + group_name
                title_box.text_frame.font.size = Pt(18)
                title_box.text_frame.font.bold = True
                top_margin += 0.5

                for student in data:
                    # if row == 0:
                    #     title_box = slide.shapes.add_textbox(Inches(left_margin),Inches(top_margin),Inches(textbox_width),Inches(textbox_height))
                    #     title_box.text_frame.text = file_path.split('/')[-1].split('_')[-1][:-5]
                    #     top_margin += 0.5
                    #     # slide = presentation.slides.add_slide(blank_slide_layout)
                    # if row > 0:
                    # if number_of_students_on_slide == 36:
                    #     slide = presentation.slides.add_slide(blank_slide_layout)
                    #     number_of_students_on_slide = 0
                    #     left_margin = 0.5
                    #     top_margin = 1
                    if number_of_students_on_slide == 12 or number_of_students_on_slide == 24:
                        left_margin += 3
                        top_margin = 1
                    new_textbox = slide.shapes.add_textbox(Inches(left_margin),Inches(top_margin),Inches(textbox_width),Inches(textbox_height))
                    self.add_student_to_slide(slide, student, new_textbox)
                    number_of_students_on_slide += 1
                    top_margin += 0.5
                    row += 1
        file_path = '{}/PowerPoints/test_advisory.pptx'.format(get_current_directory())
        presentation.save(file_path)
        print('- Saved PowerPoint at ./{}'.format('/'.join(file_path.split('/')[-2:])))
