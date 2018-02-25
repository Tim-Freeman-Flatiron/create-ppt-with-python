from pptx import Presentation
from pptx.util import Inches

def make_all_ppts(spreadsheet_data):
  for data_object in spreadsheet_data:
    file_path = data_object['file']
    student_data = data_object['student_data']
    make_ppt(file_path, student_data)

def add_student_to_slide(slide, student, textbox):
  first_name = student[0]
  last_name = student[1]
  data_item = student[2]
  textbox.text_frame.text = "{} {} {}".format(first_name, last_name, data_item)

def make_ppt(file_path, student_data):
  print('Making PowerPoint...')
  presentation = Presentation()
  blank_slide_layout = presentation.slide_layouts[6]
  slide = presentation.slides.add_slide(blank_slide_layout)

  row = 0
  number_of_students_on_slide = 0
  left_margin = 0.5
  top_margin = 1
  textbox_width = 1
  textbox_height = 1

  for student in student_data:
    if row > 0:
      if number_of_students_on_slide == 34:
        slide = presentation.slides.add_slide(blank_slide_layout)
        number_of_students_on_slide = 0
        left_margin = 0.5
        top_margin = 1
      if number_of_students_on_slide == 12 or number_of_students_on_slide == 24:
        left_margin += 3
        top_margin = 1
      new_textbox = slide.shapes.add_textbox(Inches(left_margin),Inches(top_margin),Inches(textbox_width),Inches(textbox_height))
      add_student_to_slide(slide, student, new_textbox)
      number_of_students_on_slide += 1
      top_margin += 0.5
    row += 1

  presentation.save(file_path)
  print('- Saved PowerPoint at ./{}'.format('/'.join(file_path.split('/')[-2:])))