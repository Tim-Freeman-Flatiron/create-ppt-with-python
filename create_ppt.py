# the email creation and sending portion of this code was inspired by http://naelshiab.com/tutorial-send-email-python/
# the ppt portion of this code was inspired by https://www.youtube.com/watch?v=ia9wXLq34us

# for path manipulation
from os import getcwd as get_current_directory

# for google api calls
from apiclient import discovery as api
from httplib2 import Http
from credentials import get_credentials # relative import of function to verify credentials

# for making ppts
from pptx import Presentation
from pptx.util import Inches

# for sending email
from smtplib import SMTP as StartEmailServer
from email.mime.multipart import MIMEMultipart as MakeComplexEmailObject
from email.mime.text import MIMEText as MakeEmailBody
from email.mime.base import MIMEBase as MakeAttachmentObject
from email.encoders import encode_base64

# relative import of private variables
from secrets import FROM_EMAIL, FROM_PASSWORD, TO_EMAIL, SHEET_ID

def convert_file_to_attachment(file):
	file_name = file.split('/')[-1] # extract file name from absolute path
	file_object = open(file, 'rb')
	attachment = MakeAttachmentObject('application', 'octet-stream')
	attachment.set_payload(file_object.read())
	encode_base64(attachment)
	attachment.add_header('Content-Disposition', "attachment; filename= %s" % file_name)
	return attachment

def make_email(files_to_attach):
	print('Creating email...')
	email = MakeComplexEmailObject()
	email_body = 'See attached.'
	email.attach(MakeEmailBody(email_body, 'plain'))

	email['From'] = FROM_EMAIL
	email['To'] = TO_EMAIL
	email['Subject'] = 'This Week\'s Community Meeting Powerpoints'
	print('- Attaching PowerPoints...')
	for file in files_to_attach:
		attachment = convert_file_to_attachment(file)
		email.attach(attachment)
		print('-- Attached {}'.format(file.split('/')[-1]))

	packaged_email = email.as_string()
	print('- Email created.')
	return packaged_email

def send_email(email):
	print('Starting server...')
	server = StartEmailServer('smtp.gmail.com', 587)
	server.starttls()
	print('- Logging into Gmail...')
	server.login(FROM_EMAIL, FROM_PASSWORD)
	print('- Sending email...')
	server.sendmail(FROM_EMAIL, TO_EMAIL, email)
	server.quit()
	print('- Email sent and server ended.')

def proccess_tabs(raw_tabs, api_adapter):
	processed_tabs = []

	for tab in raw_tabs:
		tab_name = tab.get('properties', {}).get('title')
		print('- Processing {}...'.format(tab_name))
		range_name = tab_name + '!A2:C'
		tab_data = api_adapter.spreadsheets().values().get(spreadsheetId=SHEET_ID, range=range_name).execute()
		student_data = tab_data.get('values', [])
		file_path = get_current_directory() + '/PowerPoints/' + tab_name + '.pptx'
		tab_object = {
									'file': file_path,
									'student_data': student_data
									}
		processed_tabs.append(tab_object)
	print('- Finished processing tabs.')
	return processed_tabs

def make_ppts_from_processed_tabs(processed_tabs):
	for tab in processed_tabs:
		file_path = tab['file']
		student_data = tab['student_data']
		make_ppt(file_path, student_data)

def make_ppt(file_path, student_data):
	print('Making PowerPoint...')
	presentation = Presentation()
	blank_slide_layout = presentation.slide_layouts[6]
	slide = presentation.slides.add_slide(blank_slide_layout)

	r = 0
	placeholder = 0
	left = 0.5
	top = 1
	width = 1
	height = 1

	for row in student_data:
		if r > 0:
			if placeholder == 34:
				slide = presentation.slides.add_slide(blank_slide_layout)
				placeholder = 0
				left = 0.5
				top = 1
			if placeholder == 12 or placeholder == 24:
				left = left + 3
				top = 1
			first_name = row[0]
			last_name = row[1]
			data_item = row[2]
			text_box = slide.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(height))
			text_box.text_frame.text = first_name + ' ' + last_name + ' ' + data_item
			placeholder = placeholder+1
			top = top+0.5
		r = r+1

	presentation.save(file_path)
	print('- Saved PowerPoint at ./{}'.format('/'.join(file_path.split('/')[-2:])))

def isolate_file_names_for_attachments(processed_tabs):
	file_names = []
	
	for tab in processed_tabs:
		file_names.append(tab['file'])

	return file_names

def main():
  print('Connecting to GoogleSheets...')
  api_credentials = get_credentials()
  http = api_credentials.authorize(Http())
  api_url = ('https://sheets.googleapis.com/$discovery/rest?''version=v4')
  api_adapter = api.build('sheets', 'v4', http=http, discoveryServiceUrl=api_url)
  print('Extracting individual tabs from the spreadsheet...')
  raw_tabs = api_adapter.spreadsheets().get(spreadsheetId=SHEET_ID).execute().get('sheets', '')
  processed_tabs = proccess_tabs(raw_tabs, api_adapter)
  make_ppts_from_processed_tabs(processed_tabs)
  files_to_attach = isolate_file_names_for_attachments(processed_tabs)
  email = make_email(files_to_attach)
  send_email(email)
  print('Done.')

if __name__ == '__main__':
	main()